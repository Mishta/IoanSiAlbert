#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
POLARIS COMMAND — Dual PowerPoint Generator  v2.0
Outputs:
  POLARIS_2019.pptx  — Office 2019 compatible (Fade transitions)
  POLARIS_365.pptx   — Microsoft 365 (Morph transitions)

Run: py -3.14 generate_pptx.py
"""

import os
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
from PIL import Image as PILImage

# ── Design System: Celestial Command ──────────────────────────────────────────
C_BG        = RGBColor(0x0f, 0x13, 0x1f)
C_SURFACE   = RGBColor(0x1b, 0x1f, 0x2c)
C_SURFACE_H = RGBColor(0x26, 0x2a, 0x37)
C_PRIMARY   = RGBColor(0xff, 0xe3, 0xb7)   # gold
C_SECONDARY = RGBColor(0xbd, 0xf4, 0xff)   # cyan
C_TERTIARY  = RGBColor(0x91, 0xff, 0x89)   # green
C_TEXT      = RGBColor(0xdf, 0xe2, 0xf3)
C_TEXT_DIM  = RGBColor(0xd4, 0xc5, 0xab)
C_OUTLINE   = RGBColor(0x4f, 0x46, 0x32)
C_DIM       = RGBColor(0x31, 0x34, 0x42)
C_WHITE     = RGBColor(0xff, 0xff, 0xff)

F_HEAD  = "Futura Md BT"   # headlines — NASA/space aesthetic (installed)
F_BODY  = "Segoe UI"        # body text — crisp, readable
F_LABEL = "Segoe UI"        # labels

W = Inches(13.333)
H = Inches(7.5)

BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# ── Slide Definitions (ROMANA — juriul cere limba romana) ──────────────────────
SLIDES = [
    # 00 — COPERTA
    {
        "num": "00", "type": "cover",
        "title": "P.O.L.A.R.I.S.",
        "subtitle": "Platforma Orbitala Laser pt.\nAlimentare si Receptie Inovativa Solara",
        "body": "Ioan CHELARU  ·  Albert OLARIU\nONCS 2026  ·  Sectiunea A  ·  Stiinte Exacte",
        "accent": C_PRIMARY,
        "image": "images/generated/06 POLARIS System General View (v2).png",
    },
    # 01 — PROBLEMA  (criteriu a — relevanta practica)
    {
        "num": "01", "type": "full_bleed", "section": "CONTEXT",
        "title": "PROBLEMA",
        "body": (
            "Cererea globala de energie creste an dupa an.\n"
            "Combustibilii fosili sunt limitati si poluanti.\n"
            "Atmosfera blocheaza 30% din energia solara.\n\n"
            "Solutia nu se afla pe Pamant.\n"
            "Solutia se afla in spatiu."
        ),
        "accent": C_SECONDARY,
        "image": "images/generated/Image A - Earth Energy Crisis.png",
    },
    # 02 — VIZIUNEA + ACRONIM
    {
        "num": "02", "type": "billboard", "section": "VIZIUNE",
        "title": (
            "Nu este o masina.\n"
            "Este o arhitectura orbitala\n"
            "care transforma lumina stelara\n"
            "in putere planetara."
        ),
        "label": "P  ·  O  ·  L  ·  A  ·  R  ·  I  ·  S",
        "accent": C_PRIMARY,
        "image": "images/generated/Image C - Full System Panoramic.png",
    },
    # 03 — OBIECTIVELE CERCETARII  (criteriu c = 10p)
    {
        "num": "03", "type": "objectives", "section": "OBIECTIVE",
        "title": "OBIECTIVELE\nCERCETARII",
        "objectives": [
            "Demonstrarea fezabilitatii teoretice a recoltarii energiei solare la 7 milioane km fata de Soare",
            "Calculul puterii prin lantul complet: oglinzi → laser → Lagrange → YBCO → receptor polar",
            "Proiectarea arhitecturii POLARIS: roi Dyson, releu Lagrange, cabluri YBCO, statii plutitoare",
            "Evaluarea eficientei sistemului si compararea cu cererea globala de energie (IEA 2023: 17,7 TW)",
            "Identificarea limitarilor actuale si a directiilor de cercetare pentru deceniile 2030–2080",
        ],
        "accent": C_TERTIARY,
    },
    # 04 — ROIUL DYSON
    {
        "num": "04", "type": "split", "section": "ARHITECTURA · NIVEL 1",
        "title": "ROIURI\nDYSON",
        "body": (
            "30 oglinzi orbitale, fiecare cu diametrul de 1 km,\n"
            "pozitionate la 7 milioane km fata de Soare.\n\n"
            "Intensitate solara la aceasta distanta:\n"
            "I7m = 6,22 x 10^5 W/m2\n\n"
            "Putere incidenta totala:\n"
            "P_incident = 1,46 x 10^13 W"
        ),
        "accent": C_PRIMARY,
        "image": "images/generated/Image B - Mirror and Sun.png",
    },
    # 05 — AVANTAJUL PROXIMITATII
    {
        "num": "05", "type": "data_hero", "section": "FIZICA",
        "title": "AVANTAJUL\nPROXIMITATII",
        "stat": "450x",
        "stat_label": "mai multa energie fata de orbita Pamantului",
        "body": (
            "La 7 milioane km fata de Soare,\n"
            "intensitatea solara urmeaza legea inversului patratului:\n\n"
            "I(r) = P_soare / 4 * pi * r2\n\n"
            "Acesta este fundamentul fizic al intregului sistem."
        ),
        "accent": C_PRIMARY,
        "image": "images/generated/01 Mirrors (v2).png",
    },
    # 06 — CONVERSIA LASER
    {
        "num": "06", "type": "split", "section": "ARHITECTURA · NIVEL 2",
        "title": "CONVERSIA\nLASER",
        "body": (
            "4 sateliti colectori transforma\n"
            "lumina solara reflectata\n"
            "in fascicule laser de inalta putere.\n\n"
            "Eficienta de conversie: 43%\n\n"
            "Energia solara bruta devine\n"
            "un fascicul directional si controlabil."
        ),
        "accent": C_SECONDARY,
        "image": "images/generated/02 Satelit Colector (close-up) (v2).png",
    },
    # 07 — NODURILE LAGRANGE
    {
        "num": "07", "type": "split", "section": "ARHITECTURA · NIVEL 3",
        "title": "NODURILE\nLAGRANGE",
        "body": (
            "L4 si L5 — puncte gravitational stabile\n"
            "in sistemul Soare-Pamant.\n\n"
            "POLARIS le foloseste ca noduri de releu energetic:\n"
            "P_L4 = 2,83 TW per nod\n\n"
            "Eficienta receptiei la Lagrange:\n"
            "eta_centralizare = 60%"
        ),
        "accent": C_TERTIARY,
        "image": "images/generated/00 System Diagram.png",
    },
    # 08 — POZITIONAREA LAGRANGE (motivatie)
    {
        "num": "08", "type": "split", "section": "MECANICA ORBITALA",
        "title": "RATIUNEA\nPOZITIONARII",
        "body": (
            "Problema: satelitii colectori orbiteaza rapid\n"
            "in apropierea Soarelui. Fara un releu fix,\n"
            "fasciculul laser ar baleia imprevizibil Pamantul.\n\n"
            "Solutia Lagrange: L4 si L5 sunt stabile la\n"
            "60° fata de Pamant, la 1 UA fata de Soare.\n\n"
            "Colector → L4/L5 (mobil→fix)\n"
            "L4/L5 → Statie Polara (fix→fix)\n\n"
            "Fasciculul nu oscileaza deasupra zonelor locuite."
        ),
        "accent": C_SECONDARY,
        "image": "images/generated/03 Satelit Lagrange (close-up) - L4-L5  (v2).png",
        "wide_text": True,
    },
    # 09 — METODOLOGIE & COMPLEXITATE  (criteriu d = 10p)
    {
        "num": "09", "type": "methodology", "section": "METODOLOGIE",
        "title": "METODOLOGIE\n& COMPLEXITATE",
        "chain": "I(r)  →  P_incident  →  eta_oglinzi  →  eta_laser  →  P_L4  →  eta_YBCO  →  P_livrata",
        "body": (
            "6 domenii integrate: fizica  ·  astronomie  ·  inginerie  ·  supraconductivitate  ·  criogenie  ·  mecanica fluidelor\n\n"
            "Modele matematice inlantuite: oglinzi → sateliti → L4/L5 → baloane → YBCO → retea\n\n"
            "Rezultat final: 1,6 TW livrati  (2 x 800 GW)\n"
            "Surse: IEA 2023, NASA NIAC, Parker Solar Probe, CERN, JAXA"
        ),
        "accent": C_SECONDARY,
    },
    # 10 — SUPRACONDUCTORI YBCO
    {
        "num": "10", "type": "split", "section": "MATERIALE · NIVEL 4",
        "title": "SUPRACONDUCTORI\nYBCO",
        "body": (
            "Cabluri YBCO (itriu-bariu-cupru-oxid) —\n"
            "inspiratie din consulturi cu fizicieni CERN.\n\n"
            "Temperatura de operare: 77 K (azot lichid)\n"
            "vs. NbTi la CERN LHC: 1,9 K (heliu lichid)\n\n"
            "Masa per cablu: ~7 tone\n"
            "Densitate curent: j <= j_c\n"
            "Pierderi de transmisie: aproape zero"
        ),
        "accent": C_SECONDARY,
        "image": "images/generated/05 Cablul YBCO (v2).png",
    },
    # 11 — NUMERELE  (1.6 TW)
    {
        "num": "11", "type": "data_hero", "section": "DATE",
        "title": "NUMERELE",
        "stat": "1,6 TW",
        "stat_label": "livrati pe Pamant",
        "body": (
            "2 statii polare x 800 GW fiecare\n"
            "= 48% din consumul global de electricitate\n\n"
            "Cerere globala: 17,7 TW  (IEA 2023)\n"
            "Eficienta cumulata a lantului: eta = 11%"
        ),
        "accent": C_TERTIARY,
        "image": "images/generated/Image D - Energy Mathematics.png",
    },
    # 12 — STATIILE PLUTITOARE
    {
        "num": "12", "type": "split", "section": "ARHITECTURA · NIVEL 5",
        "title": "STATII\nPLUTITOARE",
        "body": (
            "Fiecare statie polara cantareste 4 tone\n"
            "si este suspendata la 30 km altitudine\n"
            "prin baloane cu heliu.\n\n"
            "Principiul lui Arhimede —\n"
            "aplicat la scara planetara.\n\n"
            "Pozitie: deasupra polilor magnetici,\n"
            "zona cu cea mai mica densitate de populatie\n"
            "din calea fasciculului laser."
        ),
        "accent": C_PRIMARY,
        "image": "images/generated/04 Polar Station (v2).png",
    },
    # 13 — STAR POWER GRID + VIITOR  (criteriu f — dezvoltare ulterioara)
    {
        "num": "13", "type": "split", "section": "EXTENSIE · VIITOR",
        "title": "STAR POWER\nGRID",
        "body": (
            "Energia surplus de la L4/L5 nu se opreste acolo.\n\n"
            "Reteaua Star Power Grid extinde transmisia\n"
            "prin microunde GEO (orbita geostationa)\n"
            "catre regiuni indepartate si slab deservite.\n\n"
            "Energie curata, fara infrastructura terestra.\n\n"
            "Orizont temporal: 30 – 50 ani"
        ),
        "accent": C_SECONDARY,
        "image": "images/generated/Image E - Earth Grid Distribution.png",
    },
    # 14 — SCALA KARDASHEV
    {
        "num": "14", "type": "split", "section": "CONTEXT CIVILIZATIONAL",
        "title": "SCALA\nKARDASHEV",
        "body": (
            "Tip 0: energie din materie organica moarta\n"
            "Tip I: toata energia disponibila pe planeta\n"
            "Tip II: toata energia unui sistem solar\n\n"
            "Omenirea se afla azi la ~0,73 pe scala.\n\n"
            "POLARIS este puntea spre Tipul I —\n"
            "nu doar un proiect energetic,\n"
            "ci o tranzitie civilizationala."
        ),
        "accent": C_PRIMARY,
        "image": "images/generated/Image F - Kardashev Scale.png",
    },
    # 15 — ARHITECTURA ORIGINALA  (criteriu a — creativitate)
    {
        "num": "15", "type": "split", "section": "INOVATIE",
        "title": "ARHITECTURA\nORIGINALA",
        "body": (
            "Nicio lucrare existenta nu combina:\n\n"
            "·  Roi Dyson la 7 milioane km fata de Soare\n"
            "·  Noduri de releu energetic Lagrange L4/L5\n"
            "·  Cabluri supraconductoare HTS YBCO\n"
            "·  Statii polare suspendate balonistic\n\n"
            "POLARIS este o arhitectura cu adevarat noua."
        ),
        "accent": C_TERTIARY,
        "image": "images/generated/06 POLARIS System General View (v2).png",
    },
    # 16 — STIINTA REALA  (criteriu b — abordare stiintifica)
    {
        "num": "16", "type": "split", "section": "VALIDARE",
        "title": "FUNDAMENT\nSTIINTIFIC",
        "body": (
            "·  NASA Parker Solar Probe\n"
            "   → a ajuns la 6,1 milioane km fata de Soare\n\n"
            "·  CERN — cercetare YBCO & consultare directa\n\n"
            "·  JAXA — transmisie wireless de energie\n"
            "   prin microunde\n\n"
            "·  NASA NIAC — studii de fezabilitate SBSP\n\n"
            "POLARIS se construieste pe stiinta de varf."
        ),
        "accent": C_SECONDARY,
        "image": "images/generated/Image G - Science Lab meets Space.png",
    },
    # 17 — LIMITARI  (maturitate stiintifica)
    {
        "num": "17", "type": "center", "section": "MATURITATE STIINTIFICA",
        "title": "LIMITARI\nCUNOSCUTE",
        "body": (
            "·  Fabricarea si lansarea oglinzilor de 1 km la 7 milioane km\n"
            "   — tehnologie inexistenta in prezent\n\n"
            "·  Mententa robotica la distante extreme — nerezolvata\n\n"
            "·  Eficienta cumulata eta = 11% — necesita verificare experimentala\n\n"
            "·  Impactul fasciculelor laser de mare putere asupra atmosferei — partial neevaluat\n\n"
            "Recunoasterea limitarilor este parte integranta a abordarii stiintifice rigoroase."
        ),
        "accent": C_OUTLINE,
    },
    # 18 — MISIUNEA  (concluzie)
    {
        "num": "18", "type": "statement", "section": "CONCLUZIE",
        "title": (
            "POLARIS nu este doar un concept.\n"
            "Este un model teoretic coerent care demonstreaza\n"
            "ca energia solara spatiala ar putea deveni\n"
            "coloana vertebrala a infrastructurii energetice\n"
            "a umanitatii in secolul XXI."
        ),
        "accent": C_PRIMARY,
    },
    # 19 — ECHIPA + BIBLIOGRAFIE  (criteriu f — activitate echipa)
    {
        "num": "19", "type": "credits", "section": "ECHIPA",
        "title": "ECHIPA &\nBIBLIOGRAFIE",
        "accent": C_SECONDARY,
        "image": "assets/team-hero.jpeg",
    },
]


# ═══════════════════════════════════════════════════════════════════════════════
# HELPER FUNCTIONS
# ═══════════════════════════════════════════════════════════════════════════════

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_name=F_BODY, font_size=12, bold=False,
             color=C_TEXT_DIM, align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_body(slide, text, left, top, width, height,
             font_name=F_BODY, font_size=18, color=C_TEXT_DIM,
             line_spacing_pct=150):
    lines = text.split("\n")
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        pPr = p._p.get_or_add_pPr()
        lnSpc = etree.SubElement(pPr, qn("a:lnSpc"))
        spcPct = etree.SubElement(lnSpc, qn("a:spcPct"))
        spcPct.set("val", str(line_spacing_pct * 1000))
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return txBox


def add_title(slide, text, left, top, width, height,
              font_size=56, color=C_PRIMARY, font_name=F_HEAD,
              align=PP_ALIGN.LEFT):
    lines = text.split("\n")
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = True
        run.font.color.rgb = color
    return txBox


def compress_image(full_path, max_px=1920, quality=82):
    """Compress image, return (BytesIO buffer, (width_px, height_px))."""
    img = PILImage.open(full_path).convert("RGB")
    w, h = img.size
    if w > max_px or h > max_px:
        ratio = min(max_px / w, max_px / h)
        img = img.resize((int(w * ratio), int(h * ratio)), PILImage.LANCZOS)
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=quality, optimize=True)
    buf.seek(0)
    return buf, img.size


def try_add_image(slide, rel_path, left, top, max_w, max_h, align_center=True):
    """Add image preserving aspect ratio within bounding box."""
    full = os.path.join(BASE_DIR, rel_path)
    if not os.path.exists(full):
        print(f"  WARNING: Image not found: {full}")
        return False
    try:
        img_stream, (pw, ph) = compress_image(full)
        aspect = pw / ph
        box_aspect = max_w / max_h

        if aspect >= box_aspect:
            # Image wider than box: constrain by width
            actual_w = max_w
            actual_h = int(max_w / aspect)
        else:
            # Image taller than box: constrain by height
            actual_h = max_h
            actual_w = int(max_h * aspect)

        if align_center:
            off_l = (max_w - actual_w) // 2
            off_t = (max_h - actual_h) // 2
        else:
            off_l = off_t = 0

        slide.shapes.add_picture(img_stream, left + off_l, top + off_t, actual_w, actual_h)
        return True
    except Exception as e:
        print(f"  WARNING: Could not add image {rel_path}: {e}")
        return False


def add_full_bleed_image(slide, rel_path, opacity_pct=72):
    """Place image filling entire slide (cover/crop), then dark overlay."""
    full = os.path.join(BASE_DIR, rel_path)
    if not os.path.exists(full):
        return False
    try:
        img_stream, (pw, ph) = compress_image(full, max_px=1920, quality=80)
        img_ar = pw / ph
        slide_ar = W / H

        if img_ar > slide_ar:
            # Image wider: fit by height, crop sides
            actual_h = H
            actual_w = int(H * img_ar)
        else:
            # Image taller: fit by width, crop top/bottom
            actual_w = W
            actual_h = int(W / img_ar)

        off_l = -(actual_w - W) // 2
        off_t = -(actual_h - H) // 2

        pic = slide.shapes.add_picture(img_stream, off_l, off_t, actual_w, actual_h)
        slide.shapes._spTree.remove(pic._element)
        slide.shapes._spTree.insert(2, pic._element)

        # Dark overlay
        overlay = add_rect(slide, 0, 0, W, H, C_BG)
        solidFill = overlay.fill._xPr.find(qn("a:solidFill"))
        if solidFill is not None:
            srgbClr = solidFill.find(qn("a:srgbClr"))
            if srgbClr is not None:
                alpha = etree.SubElement(srgbClr, qn("a:alpha"))
                alpha.set("val", str(opacity_pct * 1000))
        return True
    except Exception as e:
        print(f"  WARNING: Could not add bleed image {rel_path}: {e}")
        return False


def add_morph_transition(slide):
    transition_xml = (
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'spd="slow"><p:morph origin="object"/></p:transition>'
    )
    slide._element.append(etree.fromstring(transition_xml))


def add_fade_transition(slide):
    transition_xml = (
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        'spd="slow"><p:fade/></p:transition>'
    )
    slide._element.append(etree.fromstring(transition_xml))


def add_chrome(slide, data, dark_mode=False):
    """Standard chrome: accent bar, section label, branding, footer, slide number."""
    accent = data.get("accent", C_PRIMARY)
    num    = data.get("num", "00")

    # Left accent bar
    add_rect(slide, Inches(0.42), Inches(1.0), Inches(0.05), Inches(4.5), accent)

    # Section label
    section = data.get("section", "")
    if section:
        label_color = accent if not dark_mode else C_TEXT_DIM
        add_text(slide, section,
                 Inches(0.62), Inches(0.5), Inches(7), Inches(0.3),
                 font_name=F_LABEL, font_size=9, bold=True, color=label_color)

    # POLARIS BEARS top right
    add_text(slide, "POLARIS BEARS",
             Inches(10.0), Inches(0.35), Inches(3.1), Inches(0.35),
             font_name=F_HEAD, font_size=9, bold=True,
             color=C_PRIMARY, align=PP_ALIGN.RIGHT)

    # Bottom rule
    add_rect(slide, Inches(0.42), Inches(7.05), Inches(12.5), Inches(0.02), C_OUTLINE)

    # Footer
    add_text(slide,
             "ONCS 2026  —  P.O.L.A.R.I.S.  —  ENERGIE SOLARA ORBITALA",
             Inches(0.55), Inches(7.15), Inches(10.5), Inches(0.26),
             font_name=F_LABEL, font_size=7, color=C_OUTLINE)

    # Slide number
    add_text(slide, num,
             Inches(12.2), Inches(7.15), Inches(0.9), Inches(0.26),
             font_name=F_LABEL, font_size=7, bold=True,
             color=C_OUTLINE, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# LAYOUT BUILDERS
# ═══════════════════════════════════════════════════════════════════════════════

def build_cover(prs, data):
    """Full-bleed background image, dark overlay, centered epic title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    # Background image (full bleed, dimmed)
    img_path = os.path.join(BASE_DIR, data["image"])
    if os.path.exists(img_path):
        add_full_bleed_image(slide, data["image"], opacity_pct=75)

    # Top laser line
    add_rect(slide, Inches(1.5), Inches(1.75), Inches(5.5), Inches(0.03), C_SECONDARY)

    # Main title — enormous
    add_title(slide, data["title"],
              Inches(1.5), Inches(1.85), Inches(10.0), Inches(2.4),
              font_size=100, color=C_PRIMARY)

    # Subtitle
    add_body(slide, data["subtitle"],
             Inches(1.5), Inches(3.85), Inches(9), Inches(1.3),
             font_name=F_HEAD, font_size=22, color=C_SECONDARY,
             line_spacing_pct=140)

    # Bottom rule
    add_rect(slide, Inches(1.5), Inches(5.55), Inches(10), Inches(0.025), C_OUTLINE)

    # Team + event
    add_body(slide, data["body"],
             Inches(1.5), Inches(5.7), Inches(9), Inches(0.9),
             font_name=F_LABEL, font_size=13, color=C_TEXT_DIM,
             line_spacing_pct=145)

    # POLARIS BEARS top right
    add_text(slide, "POLARIS BEARS",
             Inches(10.0), Inches(0.35), Inches(3.1), Inches(0.35),
             font_name=F_HEAD, font_size=9, bold=True,
             color=C_PRIMARY, align=PP_ALIGN.RIGHT)

    # Slide 00 counter
    add_text(slide, "00",
             Inches(12.2), Inches(7.15), Inches(0.9), Inches(0.26),
             font_name=F_LABEL, font_size=7, bold=True,
             color=C_OUTLINE, align=PP_ALIGN.RIGHT)

    return slide


def build_full_bleed(prs, data):
    """Full-bleed image as background, text overlaid on left."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    if data.get("image"):
        add_full_bleed_image(slide, data["image"], opacity_pct=70)

    add_chrome(slide, data)

    accent = data.get("accent", C_PRIMARY)

    # Large title
    add_title(slide, data["title"],
              Inches(0.62), Inches(1.1), Inches(7.5), Inches(2.0),
              font_size=72, color=accent)

    # Body
    body = data.get("body", "")
    if body:
        add_body(slide, body,
                 Inches(0.62), Inches(3.3), Inches(7.0), Inches(3.4),
                 font_size=20, color=C_TEXT, line_spacing_pct=160)

    return slide


def build_billboard(prs, data):
    """Giant quote statement. Background image. Minimal chrome."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    if data.get("image"):
        add_full_bleed_image(slide, data["image"], opacity_pct=65)

    accent = data.get("accent", C_PRIMARY)

    # POLARIS BEARS top right
    add_text(slide, "POLARIS BEARS",
             Inches(10.0), Inches(0.35), Inches(3.1), Inches(0.35),
             font_name=F_HEAD, font_size=9, bold=True,
             color=C_PRIMARY, align=PP_ALIGN.RIGHT)

    # Section label
    section = data.get("section", "")
    if section:
        add_text(slide, section,
                 Inches(0.62), Inches(0.5), Inches(7), Inches(0.3),
                 font_name=F_LABEL, font_size=9, bold=True, color=accent)

    # Left accent bar (tall)
    add_rect(slide, Inches(0.42), Inches(1.0), Inches(0.06), Inches(5.5), accent)

    # Big statement
    add_title(slide, data["title"],
              Inches(0.75), Inches(1.1), Inches(11.8), Inches(5.2),
              font_size=42, color=C_WHITE, font_name=F_HEAD,
              align=PP_ALIGN.LEFT)

    # Label at bottom
    label = data.get("label", "")
    if label:
        add_text(slide, label,
                 Inches(0.75), Inches(6.5), Inches(11), Inches(0.5),
                 font_name=F_LABEL, font_size=13, bold=True,
                 color=accent, align=PP_ALIGN.LEFT)

    # Bottom rule + footer
    add_rect(slide, Inches(0.42), Inches(7.05), Inches(12.5), Inches(0.02), C_OUTLINE)
    add_text(slide, "ONCS 2026  —  P.O.L.A.R.I.S.  —  ENERGIE SOLARA ORBITALA",
             Inches(0.55), Inches(7.15), Inches(10.5), Inches(0.26),
             font_name=F_LABEL, font_size=7, color=C_OUTLINE)
    add_text(slide, data.get("num", "02"),
             Inches(12.2), Inches(7.15), Inches(0.9), Inches(0.26),
             font_name=F_LABEL, font_size=7, bold=True,
             color=C_OUTLINE, align=PP_ALIGN.RIGHT)

    return slide


def build_objectives(prs, data):
    """Numbered objectives list. Dark background, no image. Criteriu c = 10p."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_chrome(slide, data)

    accent = data.get("accent", C_TERTIARY)

    # Title
    add_title(slide, data["title"],
              Inches(0.62), Inches(1.1), Inches(4.5), Inches(2.2),
              font_size=44, color=accent)

    # Thin vertical separator
    add_rect(slide, Inches(5.1), Inches(1.0), Inches(0.025), Inches(5.8), C_OUTLINE)

    # Objectives list
    objectives = data.get("objectives", [])
    obj_top = Inches(1.05)
    obj_step = Inches(1.12)
    for i, obj_text in enumerate(objectives):
        # Number badge background
        add_rect(slide, Inches(5.35), obj_top + Inches(0.05),
                 Inches(0.55), Inches(0.55), accent)
        # Number
        add_text(slide, str(i + 1),
                 Inches(5.35), obj_top + Inches(0.05),
                 Inches(0.55), Inches(0.55),
                 font_name=F_HEAD, font_size=16, bold=True,
                 color=C_BG, align=PP_ALIGN.CENTER)
        # Objective text
        add_text(slide, obj_text,
                 Inches(6.1), obj_top, Inches(6.9), Inches(0.95),
                 font_name=F_BODY, font_size=17, bold=False,
                 color=C_TEXT_DIM, align=PP_ALIGN.LEFT)
        obj_top += obj_step

    return slide


def build_split(prs, data):
    """Text left, image right. Most common content slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_chrome(slide, data)

    accent    = data.get("accent", C_PRIMARY)
    wide_text = data.get("wide_text", False)

    if wide_text:
        text_w  = Inches(7.8)
        img_l   = Inches(8.6)
        img_w   = Inches(4.3)
        title_sz = 40
    else:
        text_w  = Inches(6.2)
        img_l   = Inches(7.0)
        img_w   = Inches(5.9)
        title_sz = 52

    img_top = Inches(1.0)
    img_h   = Inches(5.8)

    # Title
    add_title(slide, data["title"],
              Inches(0.62), Inches(1.1), text_w, Inches(2.2),
              font_size=title_sz, color=accent)

    # Stat (data slides)
    if data.get("stat"):
        add_text(slide, data["stat"],
                 Inches(0.62), Inches(3.25), text_w, Inches(1.1),
                 font_name=F_HEAD, font_size=56, bold=True, color=accent)
        add_text(slide, data.get("stat_label", ""),
                 Inches(0.62), Inches(4.25), text_w, Inches(0.4),
                 font_name=F_LABEL, font_size=10, bold=True, color=C_TEXT_DIM)

    # Body
    body = data.get("body", "")
    if body:
        body_top = Inches(4.6) if data.get("stat") else Inches(3.45)
        body_sz  = 16 if wide_text else 18
        add_body(slide, body,
                 Inches(0.62), body_top, text_w, Inches(2.9),
                 font_size=body_sz, color=C_TEXT_DIM, line_spacing_pct=155)

    # Image
    if data.get("image"):
        try_add_image(slide, data["image"], img_l, img_top, img_w, img_h)

    return slide


def build_data_hero(prs, data):
    """Giant stat number + background image + supporting text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    # Background image (dimmed)
    if data.get("image"):
        add_full_bleed_image(slide, data["image"], opacity_pct=82)

    add_chrome(slide, data)

    accent = data.get("accent", C_PRIMARY)

    # Title (smaller, top)
    add_title(slide, data["title"],
              Inches(0.62), Inches(1.1), Inches(8), Inches(1.6),
              font_size=44, color=C_TEXT_DIM)

    # GIANT stat
    add_text(slide, data.get("stat", ""),
             Inches(0.62), Inches(2.5), Inches(12.0), Inches(2.2),
             font_name=F_HEAD, font_size=110, bold=True,
             color=accent, align=PP_ALIGN.LEFT)

    # Stat label
    add_text(slide, data.get("stat_label", ""),
             Inches(0.62), Inches(4.55), Inches(10), Inches(0.4),
             font_name=F_LABEL, font_size=13, bold=True,
             color=C_TEXT_DIM, align=PP_ALIGN.LEFT)

    # Horizontal rule
    add_rect(slide, Inches(0.62), Inches(5.1), Inches(11.5), Inches(0.02), C_OUTLINE)

    # Supporting body
    body = data.get("body", "")
    if body:
        add_body(slide, body,
                 Inches(0.62), Inches(5.2), Inches(11.5), Inches(1.6),
                 font_size=17, color=C_TEXT_DIM, line_spacing_pct=145)

    return slide


def build_methodology(prs, data):
    """Methodology slide: chain equation displayed + 6-domain summary. Criteriu d = 10p."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_chrome(slide, data)

    accent = data.get("accent", C_SECONDARY)

    # Title
    add_title(slide, data["title"],
              Inches(0.62), Inches(1.1), Inches(7), Inches(1.9),
              font_size=48, color=accent)

    # Chain equation — prominent box
    chain = data.get("chain", "")
    if chain:
        # Dark surface box behind equation
        add_rect(slide, Inches(0.42), Inches(3.1), Inches(12.5), Inches(0.85), C_SURFACE)
        add_text(slide, chain,
                 Inches(0.6), Inches(3.2), Inches(12.3), Inches(0.65),
                 font_name="Consolas", font_size=16, bold=True,
                 color=C_PRIMARY, align=PP_ALIGN.CENTER)

    # Body
    body = data.get("body", "")
    if body:
        add_body(slide, body,
                 Inches(0.62), Inches(4.15), Inches(12.3), Inches(2.6),
                 font_size=18, color=C_TEXT_DIM, line_spacing_pct=155)

    return slide


def build_center(prs, data):
    """Text-only, no image. Wide layout."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)
    add_chrome(slide, data)

    accent = data.get("accent", C_PRIMARY)

    add_title(slide, data["title"],
              Inches(0.62), Inches(1.1), Inches(12.0), Inches(2.0),
              font_size=52, color=accent)

    body = data.get("body", "")
    if body:
        add_body(slide, body,
                 Inches(0.62), Inches(3.3), Inches(12.0), Inches(3.5),
                 font_size=19, color=C_TEXT_DIM, line_spacing_pct=155)

    return slide


def build_statement(prs, data):
    """Full-screen manifesto — no image, big centered text."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    accent = data.get("accent", C_PRIMARY)

    # Large left accent line
    add_rect(slide, Inches(0.42), Inches(0.8), Inches(0.08), Inches(5.9), accent)

    # POLARIS BEARS
    add_text(slide, "POLARIS BEARS",
             Inches(10.0), Inches(0.35), Inches(3.1), Inches(0.35),
             font_name=F_HEAD, font_size=9, bold=True,
             color=C_PRIMARY, align=PP_ALIGN.RIGHT)

    section = data.get("section", "")
    if section:
        add_text(slide, section,
                 Inches(0.62), Inches(0.5), Inches(7), Inches(0.3),
                 font_name=F_LABEL, font_size=9, bold=True, color=accent)

    # Giant statement text
    add_title(slide, data["title"],
              Inches(0.75), Inches(0.95), Inches(12.0), Inches(5.8),
              font_size=34, color=C_WHITE, font_name=F_HEAD,
              align=PP_ALIGN.LEFT)

    # Bottom chrome
    add_rect(slide, Inches(0.42), Inches(7.05), Inches(12.5), Inches(0.02), C_OUTLINE)
    add_text(slide, "ONCS 2026  —  P.O.L.A.R.I.S.  —  ENERGIE SOLARA ORBITALA",
             Inches(0.55), Inches(7.15), Inches(10.5), Inches(0.26),
             font_name=F_LABEL, font_size=7, color=C_OUTLINE)
    add_text(slide, data.get("num", "18"),
             Inches(12.2), Inches(7.15), Inches(0.9), Inches(0.26),
             font_name=F_LABEL, font_size=7, bold=True,
             color=C_OUTLINE, align=PP_ALIGN.RIGHT)

    return slide


def build_credits(prs, data):
    """Team + bibliography. Photo on right."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, C_BG)

    accent = data.get("accent", C_SECONDARY)

    # Left accent bar
    add_rect(slide, Inches(0.42), Inches(0.5), Inches(0.06), Inches(6.6), accent)

    # Section
    add_text(slide, "CREDITS",
             Inches(0.62), Inches(0.4), Inches(3), Inches(0.3),
             font_name=F_LABEL, font_size=9, bold=True, color=accent)

    # POLARIS BEARS
    add_text(slide, "POLARIS BEARS",
             Inches(10.0), Inches(0.35), Inches(3.1), Inches(0.35),
             font_name=F_HEAD, font_size=9, bold=True,
             color=C_PRIMARY, align=PP_ALIGN.RIGHT)

    # Title
    add_title(slide, data["title"],
              Inches(0.62), Inches(0.65), Inches(6.5), Inches(1.5),
              font_size=38, color=accent)

    # Team column header
    add_text(slide, "ECHIPA",
             Inches(0.62), Inches(2.25), Inches(5), Inches(0.35),
             font_name=F_LABEL, font_size=9, bold=True, color=C_PRIMARY)
    add_rect(slide, Inches(0.62), Inches(2.6), Inches(5.8), Inches(0.02), C_OUTLINE)

    team_text = (
        "Ioan CHELARU\n"
        "  Arhitectura sistem, calcule energetice,\n"
        "  supraconductori YBCO, dinamica orbitala\n\n"
        "Albert OLARIU\n"
        "  Modelare vizuala, scala Kardashev,\n"
        "  reteaua Star Power Grid, prezentare\n\n"
        "Colaboratori:\n"
        "  — de completat —"
    )
    add_body(slide, team_text,
             Inches(0.62), Inches(2.72), Inches(5.9), Inches(4.0),
             font_size=14, color=C_TEXT_DIM, line_spacing_pct=140)

    # Team photo
    if data.get("image"):
        try_add_image(slide, data["image"],
                      Inches(6.8), Inches(1.0), Inches(3.8), Inches(5.0))

    # Bibliography
    add_text(slide, "BIBLIOGRAFIE — STIL HARVARD",
             Inches(0.62), Inches(4.95),
             Inches(12), Inches(0.3),
             font_name=F_LABEL, font_size=8, bold=True, color=C_SECONDARY)

    bib = (
        "Mankins (2011) · ESA (2023) · Glaser (1968) · Benford & Benford (2015) · "
        "Kardashev (1964) · NASA NIAC · Parker Solar Probe · Lagrange Points · JAXA SSPS · "
        "IEA World Energy Outlook 2023 · Chou et al. (2019) · Aschenbrenner (2023)"
    )
    add_text(slide, bib,
             Inches(0.62), Inches(5.35), Inches(12.5), Inches(0.8),
             font_name=F_LABEL, font_size=8, color=C_OUTLINE)

    # Bottom chrome
    add_rect(slide, Inches(0.42), Inches(7.05), Inches(12.5), Inches(0.02), C_OUTLINE)
    add_text(slide, "ONCS 2026  —  P.O.L.A.R.I.S.  —  ENERGIE SOLARA ORBITALA",
             Inches(0.55), Inches(7.15), Inches(10.5), Inches(0.26),
             font_name=F_LABEL, font_size=7, color=C_OUTLINE)
    add_text(slide, "19",
             Inches(12.2), Inches(7.15), Inches(0.9), Inches(0.26),
             font_name=F_LABEL, font_size=7, bold=True,
             color=C_OUTLINE, align=PP_ALIGN.RIGHT)

    return slide


# ═══════════════════════════════════════════════════════════════════════════════
# DISPATCHER
# ═══════════════════════════════════════════════════════════════════════════════

def build_slide(prs, data):
    stype = data.get("type", "split")
    dispatch = {
        "cover":       build_cover,
        "full_bleed":  build_full_bleed,
        "billboard":   build_billboard,
        "objectives":  build_objectives,
        "split":       build_split,
        "data_hero":   build_data_hero,
        "methodology": build_methodology,
        "center":      build_center,
        "statement":   build_statement,
        "credits":     build_credits,
    }
    builder = dispatch.get(stype, build_split)
    return builder(prs, data)


# ═══════════════════════════════════════════════════════════════════════════════
# GENERATE
# ═══════════════════════════════════════════════════════════════════════════════

def generate(m365=False):
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H

    for data in SLIDES:
        label = data.get("title", "?")
        label = label.replace("\n", " ")[:35]
        print(f"  [{data['num']}] {label}...")
        slide = build_slide(prs, data)
        if m365:
            add_morph_transition(slide)
        else:
            add_fade_transition(slide)

    suffix   = "365" if m365 else "2019"
    out_path = os.path.join(BASE_DIR, f"POLARIS_{suffix}.pptx")
    prs.save(out_path)
    size_mb = os.path.getsize(out_path) / 1024 / 1024
    print(f"  Saved: POLARIS_{suffix}.pptx  ({size_mb:.1f} MB)")
    return out_path


if __name__ == "__main__":
    print("=" * 60)
    print("POLARIS COMMAND - PowerPoint Generator v2.0")
    print("=" * 60)
    print("\n[1/2] Generating POLARIS_2019.pptx ...")
    generate(m365=False)
    print("\n[2/2] Generating POLARIS_365.pptx ...")
    generate(m365=True)
    print("\n" + "=" * 60)
    print("Done. Deschide fisierele in PowerPoint pentru preview.")
    print("=" * 60)
